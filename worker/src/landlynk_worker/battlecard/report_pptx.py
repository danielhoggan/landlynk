"""Multi-slide report deck from a Battlecard payload.

The fuller, unbranded-by-default analytical report (the LA Insight equivalent):
a cover, an overview, the AI area profile, the three demographic charts with
narrative, a marketing development plan and a methodology appendix. It renders
from the same single Battlecard payload as the one-slide card, and themes to the
brand (primary, secondary, accent colours and logo) when those are supplied.

One area is one report (use the combined Battlecard for a whole-catchment study
area). Charts are native and editable. Generated prose already follows the house
conventions.
"""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
)
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.slide import Slide
from pptx.util import Inches, Length, Pt

from .schema import Battlecard, DataValue

# LandLynk green, the brand wordmark colour. Used as the default heading and
# theme colour when no client brand colour is supplied (see web tailwind config).
_DEFAULT_HEADING = "2F6B3A"
_FONT = "Poppins"
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_INK = RGBColor(0x1E, 0x2A, 0x32)
_GREY = RGBColor(0x6A, 0x73, 0x7B)
_LIGHT = RGBColor(0xF4, 0xF5, 0xF7)
_GOLD = RGBColor(0xC9, 0xA2, 0x4B)
_TENURE_COLORS = [
    RGBColor(0xC0, 0x4A, 0x1F),
    RGBColor(0x1F, 0x5A, 0x3C),
    RGBColor(0x0A, 0x1F, 0x44),
    RGBColor(0x8E, 0x95, 0x9C),
]
_AMENITY_ORDER = ["Transport", "Retail", "Leisure", "Education", "Healthcare", "Other"]


def _hex(color: str | None, fallback: str = _DEFAULT_HEADING) -> RGBColor:
    return RGBColor.from_string((color or fallback).lstrip("#"))


def _darken(color: RGBColor, factor: float = 0.6) -> RGBColor:
    """A darker shade of a colour, for panels on a coloured slide background."""
    return RGBColor(
        int(color[0] * factor), int(color[1] * factor), int(color[2] * factor)
    )


def _num(dv: DataValue) -> float:
    return 0.0 if dv.value is None else float(dv.value)


def _money(value: float | None) -> str:
    if value is None:
        return "n/a"
    if value >= 1_000_000:
        return f"£{value / 1_000_000:.1f}m"
    if value >= 1_000:
        return f"£{value / 1_000:.0f}k"
    return f"£{value:,.0f}"


# Each housebuilder intent frames the deck: a cover eyebrow, a one-line purpose
# (formatted with the audience), and the strategy slide title.
_INTENTS: dict[str, dict[str, str]] = {
    "find_site": {
        "eyebrow": "SITE FINDING REPORT",
        "purpose": "Where to build for {audience} across the search area",
        "plan": "Where to build",
    },
    "appraise": {
        "eyebrow": "SITE APPRAISAL",
        "purpose": "Is this site worth acquiring, and who would buy here",
        "plan": "Site appraisal summary",
    },
    "next_phase": {
        "eyebrow": "NEXT-PHASE PLAN",
        "purpose": "What to build for the remaining phase of this site",
        "plan": "Next-phase product plan",
    },
}


def render_report_pptx(
    card: Battlecard,
    heading_color: str | None = None,
    logo: bytes | None = None,
    accent: str | None = None,
    secondary: str | None = None,
    map_image: bytes | None = None,
    development_context: dict | None = None,
    intent: str | None = None,
    audience_label: str | None = None,
    supply: dict | None = None,
) -> bytes:
    """Render the full report deck for one Battlecard payload."""
    navy = _hex(heading_color)
    accent_rgb = _hex(accent, "C9A24B")
    secondary_rgb = _hex(secondary, "1F5A3C")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    framing = _INTENTS.get(intent or "")
    theme = {
        "navy": navy,
        "accent": accent_rgb,
        "secondary": secondary_rgb,
        "framing": framing,
        "audience": audience_label,
        "supply": supply,
    }
    _cover(prs.slides.add_slide(blank), card, theme, logo)
    # Content slides, each branded with a footer mark (brand logo or wordmark).
    for build in (
        lambda s: _overview(s, card, theme, development_context),
        lambda s: _area_profile(s, card, theme, map_image),
        lambda s: _age(s, card, theme),
        lambda s: _income(s, card, theme),
        lambda s: _tenure(s, card, theme),
    ):
        slide = prs.slides.add_slide(blank)
        build(slide)
        _footer(slide, theme, logo, on_dark=False)
    plan = prs.slides.add_slide(blank)
    _plan(plan, card, theme)
    _footer(plan, theme, logo, on_dark=True)
    appendix = prs.slides.add_slide(blank)
    _appendix(appendix, card, theme)
    _footer(appendix, theme, logo, on_dark=False)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- primitives -------------------------------------------------------------


def _rect(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    color: RGBColor,
) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def _text(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    lines: list[tuple],
    shrink: bool = False,
) -> None:
    """Each line is (text, size, color, bold, italic).

    With ``shrink`` the text auto-fits its box, so long generated prose (an AI
    area description, amenity lists) never overflows into the next block.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    if shrink:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, (text, size, color, bold, italic) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = _FONT
        run.font.color.rgb = color


def _section(slide: Slide, label: str, title: str, navy: RGBColor) -> None:
    _text(
        slide,
        Inches(0.6),
        Inches(0.45),
        Inches(11),
        Inches(0.3),
        [(label.upper(), 11, _GOLD, True, False)],
    )
    _text(
        slide,
        Inches(0.6),
        Inches(0.75),
        Inches(11),
        Inches(0.6),
        [(title, 26, navy, True, False)],
    )


def _wordmark(
    slide: Slide,
    left: Length,
    top: Length,
    height: Length,
    theme: dict,
    on_dark: bool,
) -> None:
    """The LandLynk wordmark, drawn as text so no image asset is needed.

    Green "Land" then "Lynk", the second half white on a dark slide and ink on
    a light one. The deck fallback when a run has no client brand logo.
    """
    box = slide.shapes.add_textbox(left, top, Inches(2.2), height)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    size = 13
    land = p.add_run()
    land.text = "Land"
    lynk = p.add_run()
    lynk.text = "Lynk"
    # On a dark cover the brand green is hard to read, so lift the whole mark to
    # the accent; on light slides keep the green "Land".
    land.font.color.rgb = theme["accent"] if on_dark else _hex(_DEFAULT_HEADING)
    lynk.font.color.rgb = _WHITE if on_dark else _INK
    for run in (land, lynk):
        run.font.name = _FONT
        run.font.bold = True
        run.font.size = Pt(size)


def _brandmark(
    slide: Slide,
    left: Length,
    top: Length,
    height: Length,
    logo: bytes | None,
    theme: dict,
    on_dark: bool,
) -> None:
    """Place the client brand logo, or fall back to the LandLynk wordmark."""
    if logo:
        try:
            slide.shapes.add_picture(io.BytesIO(logo), left, top, height=height)
            return
        except Exception:  # pragma: no cover - bad image bytes, fall back to text
            pass
    _wordmark(slide, left, top, height, theme, on_dark)


def _footer(slide: Slide, theme: dict, logo: bytes | None, on_dark: bool = False) -> None:
    """A small brand mark in the slide footer, so every page carries branding."""
    _brandmark(
        slide, Inches(11.0), Inches(6.95), Inches(0.32), logo, theme, on_dark
    )


def _bar(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    cats: list[str],
    vals: list[float],
    color: RGBColor,
    value_format: str = '0"%"',
) -> None:
    data = CategoryChartData()
    data.categories = cats
    data.add_series("", vals)
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data
    )
    chart = frame.chart
    chart.has_legend = False
    chart.has_title = False
    try:
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = color
        chart.value_axis.visible = False
        chart.value_axis.has_major_gridlines = False
        chart.category_axis.tick_labels.font.size = Pt(12)
        chart.category_axis.tick_labels.font.name = _FONT
        # Label each bar with its value, so the chart reads without a y axis.
        plot = chart.plots[0]
        plot.has_data_labels = True
        labels = plot.data_labels
        labels.number_format = value_format
        labels.number_format_is_linked = False
        labels.position = XL_LABEL_POSITION.OUTSIDE_END
        labels.font.size = Pt(12)
        labels.font.bold = True
        labels.font.name = _FONT
        labels.font.color.rgb = _INK
    except Exception:
        pass


def _donut(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    cats: list[str],
    vals: list[float],
) -> None:
    data = CategoryChartData()
    data.categories = cats
    data.add_series("", vals)
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, left, top, width, height, data
    )
    chart = frame.chart
    chart.has_title = False
    chart.has_legend = True
    try:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)
        chart.legend.font.name = _FONT
        plot = chart.plots[0]
        plot.has_data_labels = True
        labels = plot.data_labels
        labels.number_format = '0"%"'
        labels.number_format_is_linked = False
        labels.font.size = Pt(11)
        labels.font.bold = True
        labels.font.name = _FONT
        labels.font.color.rgb = _WHITE
        points = plot.series[0].points
        for i, pt in enumerate(points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = _TENURE_COLORS[i % len(_TENURE_COLORS)]
    except Exception:
        pass


# --- slides -----------------------------------------------------------------


def _cover(slide: Slide, card: Battlecard, theme: dict, logo: bytes | None) -> None:
    h = card.visual_summary.header
    framing = theme.get("framing")
    audience = theme.get("audience") or "the target buyer"
    eyebrow = (
        framing["eyebrow"] if framing else "HOUSING DEVELOPMENT INTELLIGENCE"
    )
    _rect(slide, 0, 0, Inches(13.333), Inches(7.5), theme["navy"])
    _rect(slide, 0, Inches(3.5), Inches(13.333), Inches(0.06), theme["accent"])
    _text(
        slide,
        Inches(0.9),
        Inches(1.9),
        Inches(11.5),
        Inches(0.4),
        [(eyebrow, 13, theme["accent"], True, False)],
    )
    _text(
        slide,
        Inches(0.9),
        Inches(2.4),
        Inches(11.5),
        Inches(1.0),
        [(h.development_name, 40, _WHITE, True, False)],
    )
    location = ", ".join(p for p in [h.town, h.postcode] if p)
    lines: list[tuple] = [(location or "Area analysis report", 16, _WHITE, False, False)]
    if framing:
        lines.append(
            (framing["purpose"].format(audience=audience), 13, _WHITE, False, True)
        )
    _text(slide, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.9), lines)
    _text(
        slide,
        Inches(0.9),
        Inches(6.6),
        Inches(11.5),
        Inches(0.4),
        [
            (
                "LandLynk · ONS Census 2021 and ONS income estimates",
                11,
                _GOLD,
                False,
                False,
            )
        ],
    )
    _brandmark(slide, Inches(10.6), Inches(0.55), Inches(0.7), logo, theme, True)


def _overview(
    slide: Slide, card: Battlecard, theme: dict, dev: dict | None = None
) -> None:
    ks = card.visual_summary.key_statistics
    _section(slide, "Overview", "Key statistics at a glance", theme["navy"])
    if dev and dev.get("hospital"):
        bits = [f"Nearest hospital: {dev['hospital']}"]
        if dev.get("km") is not None:
            bits.append(f"{dev['km']:g} km")
        if dev.get("ae4hr") is not None:
            bits.append(f"A&E 4-hour: {dev['ae4hr']:.0f}%")
        if dev.get("rttWeeks") is not None:
            bits.append(f"RTT median: {dev['rttWeeks']:.0f} wks")
        _text(
            slide,
            Inches(0.6),
            Inches(6.4),
            Inches(12),
            Inches(0.5),
            [("  ·  ".join(bits), 12, _INK, False, False)],
        )
    tiles = [
        (_fmt_int(ks.population_catchment), "POPULATION"),
        (_fmt_int(ks.median_age, suffix=""), "MEDIAN AGE"),
        (_money(ks.average_household_income.value), "AVG HH INCOME"),
        (_fmt_int(ks.households_catchment), "HOUSEHOLDS"),
        (_pct(ks.owner_occupied_percentage), "OWNER OCCUPIED"),
        (f"{ks.bed_range} bed", "PRODUCT"),
    ]
    for i, (value, label) in enumerate(tiles):
        col = i % 3
        row = i // 3
        left = Inches(0.6 + col * 4.05)
        top = Inches(2.0 + row * 2.0)
        _rect(slide, left, top, Inches(3.8), Inches(1.7), _LIGHT)
        _rect(slide, left, top, Inches(0.08), Inches(1.7), theme["accent"])
        _text(
            slide,
            left + Inches(0.35),
            top + Inches(0.35),
            Inches(3.3),
            Inches(0.7),
            [(value, 30, theme["navy"], True, False)],
        )
        _text(
            slide,
            left + Inches(0.35),
            top + Inches(1.1),
            Inches(3.3),
            Inches(0.4),
            [(label, 11, _GREY, True, False)],
        )


def _area_profile(
    slide: Slide, card: Battlecard, theme: dict, map_image: bytes | None
) -> None:
    _section(slide, "Section 01", "Local area profile", theme["navy"])
    profile = card.local_area_profile
    desc = profile.description if profile else card.visual_summary.header.strapline
    # Left column: the description in a bounded box (shrink-to-fit so long AI
    # prose never spills into the amenities), then amenities below with a gap.
    _text(
        slide,
        Inches(0.6),
        Inches(1.65),
        Inches(6.6),
        Inches(3.0),
        [
            ("AREA DESCRIPTION", 11, _GOLD, True, False),
            (
                desc or "Generate an AI local area profile to populate this section.",
                12,
                _INK,
                False,
                False,
            ),
        ],
        shrink=True,
    )
    if profile and profile.amenities:
        grouped: dict[str, list[str]] = {}
        for a in profile.amenities:
            grouped.setdefault(a.category, []).append(a.name)
        lines: list[tuple] = [("LOCAL AMENITIES", 11, _GOLD, True, False)]
        for cat in _AMENITY_ORDER:
            if grouped.get(cat):
                lines.append((cat.upper(), 10, theme["navy"], True, False))
                for name in grouped[cat][:3]:
                    lines.append((f"· {name}", 11, _INK, False, False))
        _text(
            slide, Inches(0.6), Inches(4.75), Inches(6.6), Inches(2.5), lines, shrink=True
        )
    # Right column: map (height-bounded so it cannot overlap the context below),
    # then local context.
    if map_image:
        try:
            slide.shapes.add_picture(
                io.BytesIO(map_image), Inches(7.5), Inches(1.65), height=Inches(3.3)
            )
        except Exception:
            pass
    if card.context_metrics:
        lines = [("LOCAL CONTEXT", 11, _GOLD, True, False)]
        for m in card.context_metrics:
            lines.append((f"{m.label}: {m.value:g} {m.unit}", 11, _INK, False, False))
        _text(
            slide, Inches(7.5), Inches(5.15), Inches(5.2), Inches(2.1), lines, shrink=True
        )


def _analysis(slide: Slide, body: str) -> None:
    """The narrative panel beside a chart. Sized to use the column it sits in."""
    _text(
        slide,
        Inches(8.2),
        Inches(1.9),
        Inches(4.6),
        Inches(4.8),
        [
            ("ANALYSIS", 13, _GOLD, True, False),
            (body, 17, _INK, False, False),
        ],
        shrink=True,
    )


def _age(slide: Slide, card: Battlecard, theme: dict) -> None:
    _section(slide, "Chart 01", "Age demographic distribution", theme["navy"])
    charts = card.visual_summary.charts
    _bar(
        slide,
        Inches(0.6),
        Inches(1.7),
        Inches(7.2),
        Inches(4.8),
        [b.label for b in charts.age_demographics],
        [_num(b.percentage) for b in charts.age_demographics],
        theme["secondary"],
    )
    cohorts = card.audience_and_demographics.age_cohorts
    body = cohorts[0].body if cohorts else ""
    _analysis(slide, body)


def _income(slide: Slide, card: Battlecard, theme: dict) -> None:
    _section(slide, "Chart 02", "Household income distribution", theme["navy"])
    inc = card.visual_summary.charts.household_income
    _bar(
        slide,
        Inches(0.6),
        Inches(1.7),
        Inches(7.2),
        Inches(4.8),
        [
            "Mean",
            "Median",
            inc.lowest_la.name or "Lowest",
            inc.highest_la.name or "Highest",
        ],
        [
            _num(inc.mean),
            _num(inc.median),
            _num(inc.lowest_la.value),
            _num(inc.highest_la.value),
        ],
        theme["secondary"],
        value_format='"£"#,##0',
    )
    _analysis(slide, card.income_and_tenure.income_commentary)


def _tenure(slide: Slide, card: Battlecard, theme: dict) -> None:
    _section(slide, "Chart 03", "Housing tenure profile", theme["navy"])
    t = card.visual_summary.charts.housing_tenure
    vals = [
        t.owns_outright.value,
        t.owns_with_mortgage.value,
        t.social_rented.value,
        t.private_rented.value,
    ]
    if any(v is not None and v > 0 for v in vals):
        _donut(
            slide,
            Inches(0.6),
            Inches(1.7),
            Inches(7.2),
            Inches(4.8),
            ["Owns outright", "Owns with mortgage", "Social rented", "Private rented"],
            [v or 0.0 for v in vals],
        )
    else:
        # No tenure data: say so rather than drawing an empty chart.
        _text(
            slide,
            Inches(0.6),
            Inches(3.2),
            Inches(7.2),
            Inches(1.6),
            [
                (
                    "Tenure data is not available for this catchment. Load the ONS "
                    "tenure dataset (TS054) to populate this chart.",
                    14,
                    _GREY,
                    False,
                    True,
                )
            ],
            shrink=True,
        )
    _analysis(slide, card.income_and_tenure.tenure_commentary)


def _plan(slide: Slide, card: Battlecard, theme: dict) -> None:
    _rect(slide, 0, 0, Inches(13.333), Inches(7.5), theme["navy"])
    _text(
        slide,
        Inches(0.6),
        Inches(0.5),
        Inches(12),
        Inches(0.4),
        [("STRATEGIC INTELLIGENCE", 11, theme["accent"], True, False)],
    )
    framing = theme.get("framing")
    title = framing["plan"] if framing else "Marketing development plan"
    _text(
        slide,
        Inches(0.6),
        Inches(0.9),
        Inches(12),
        Inches(0.6),
        [(title, 26, _WHITE, True, False)],
    )

    tiers = card.audience_and_demographics.audience_tiers
    segments = " | ".join(f"{t.tier.title()}: {t.audience}" for t in tiers) or "n/a"
    messages: list[str] = []
    for m in card.visual_summary.audience_messaging:
        messages.extend(m.message_lines)
    # Fold buildable supply and competitor schemes into the demand line when the
    # land layers are loaded, so the plan reflects the land journey.
    demand = _segments_line(card)
    supply = theme.get("supply")
    if supply and (supply.get("buildablePlots") or supply.get("competitorSchemes")):
        demand += (
            f"  •  Buildable land {supply['buildablePlots']:,} plots "
            f"({supply['buildableHomes']:,} homes)  •  Competitor schemes "
            f"{supply['competitorSchemes']:,} ({supply['competitorHomes']:,} homes)"
        )
    cards = [
        ("PRIMARY TARGET SEGMENTS", segments),
        (
            "PRODUCT AND PRICING",
            f"{card.visual_summary.key_statistics.bed_range} bed. "
            + card.pricing_rationale.positioning,
        ),
        ("KEY MARKETING MESSAGES", "  •  ".join(messages[:5]) or "n/a"),
        ("ADDRESSABLE DEMAND AND LAND", demand),
    ]
    panel = _darken(theme["navy"], 0.6)
    for i, (label, body) in enumerate(cards):
        top = Inches(1.9 + i * 1.3)
        _rect(
            slide,
            Inches(0.6),
            top,
            Inches(12.1),
            Inches(1.15),
            panel,
        )
        _rect(slide, Inches(0.6), top, Inches(0.08), Inches(1.15), theme["accent"])
        _text(
            slide,
            Inches(0.9),
            top + Inches(0.12),
            Inches(11.6),
            Inches(0.95),
            [
                (label, 11, theme["accent"], True, False),
                (body, 12, _WHITE, False, False),
            ],
        )


def _appendix(slide: Slide, card: Battlecard, theme: dict) -> None:
    _section(slide, "Appendix", "Data sources and methodology", theme["navy"])
    sources = [
        (
            "ONS Census 2021 — Age (TS007)",
            "Single year of age counts, for population, median age and cohort shares.",
        ),
        (
            "ONS Census 2021 — Household composition (TS003) and tenure (TS054)",
            "Household counts and tenure mix used for the family and tenure signals.",
        ),
        (
            "ONS income estimates for small areas",
            "Net annual household income, for the income fit and affordability.",
        ),
        (
            "ONS median house prices by MSOA",
            "Latest median price paid, for the local price context and pricing.",
        ),
    ]
    lines: list[tuple] = []
    for title, body in sources:
        lines.append((title, 13, theme["navy"], True, False))
        lines.append((body, 11, _GREY, False, False))
    dc = card.data_confidence
    lines.append((f"Data confidence: {dc.level}. {dc.note}", 11, _INK, False, True))
    _text(slide, Inches(0.6), Inches(1.7), Inches(12), Inches(5.2), lines)


# --- helpers ----------------------------------------------------------------


def _fmt_int(dv: DataValue, suffix: str = "") -> str:
    if dv.value is None:
        return "n/a"
    return f"{dv.value:,.0f}{suffix}"


def _pct(dv: DataValue) -> str:
    return "n/a" if dv.value is None else f"{dv.value:.0f}%"


def _segments_line(card: Battlecard) -> str:
    seg = card.addressable_segments

    def n(dv: DataValue) -> str:
        return "n/a" if dv.value is None else f"{dv.value:,.0f}"

    return (
        f"FTB pipeline {n(seg.first_time_buyer_pipeline)}  •  "
        f"Downsizer pool {n(seg.downsizer_pool)}  •  "
        f"Family households {n(seg.family_households)}"
    )
