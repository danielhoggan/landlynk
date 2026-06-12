"""PPTX export of a Battlecard.

The fourth render surface for the single Battlecard payload, for sales-enablement
decks. One area is one slide (the Abbots Vale reference layout): a brand sidebar
with key statistics, two messaging columns, and three charts. Brand theming is
per client and passed in, never hard-coded (design-framework.md): the heading
colour defaults to Royal Blue and the client primary overrides it. Generated
prose already follows the house conventions.
"""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide, SlideLayout
from pptx.text.text import TextFrame
from pptx.util import Inches, Length, Pt

from .schema import Battlecard, DataValue

_DEFAULT_HEADING = "4169E1"  # Royal Blue, no leading hash for python-pptx
_FONT = "Poppins"  # named only; PowerPoint substitutes if the font is not installed

_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_GOLD = RGBColor(0xC9, 0xA2, 0x4B)
_GREY = RGBColor(0x9A, 0xA3, 0xAD)
_INK = RGBColor(0x1E, 0x2A, 0x32)
_LIGHT = RGBColor(0xF1, 0xF3, 0xF5)
# Categorical, accessible palette shared with the web tenure donut.
_TENURE_COLORS = [
    RGBColor(0xC0, 0x4A, 0x1F),
    RGBColor(0x1F, 0x5A, 0x3C),
    RGBColor(0x0A, 0x1F, 0x44),
    RGBColor(0x8E, 0x95, 0x9C),
]


def _hex(color: str | None) -> RGBColor:
    value = (color or _DEFAULT_HEADING).lstrip("#")
    return RGBColor.from_string(value)


def _num(dv: DataValue) -> float:
    return 0.0 if dv.value is None else float(dv.value)


def _fmt(dv: DataValue, *, money: bool = False, pct: bool = False) -> str:
    if dv.value is None:
        return "Suppressed" if dv.suppressed else "n/a"
    if money:
        return _short_money(dv.value)
    if pct:
        return f"{dv.value:.0f}%"
    return f"{dv.value:,.0f}"


def _short_money(value: float) -> str:
    if value >= 1_000_000:
        return f"£{value / 1_000_000:.1f}m"
    if value >= 1_000:
        return f"£{value / 1_000:.0f}k"
    return f"£{value:,.0f}"


def render_battlecard_pptx(
    card: Battlecard,
    heading_color: str | None = None,
    map_geometry: dict | None = None,
    logo: bytes | None = None,
    accent: str | None = None,
    map_image: bytes | None = None,
) -> bytes:
    """Render a single Battlecard payload to PPTX bytes (one slide)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _add_card_slide(
        prs,
        prs.slide_layouts[6],
        card,
        _hex(heading_color),
        map_geometry,
        logo,
        accent,
        map_image,
    )
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def render_battlecards_pptx(
    cards: list[Battlecard],
    heading_color: str | None = None,
    logo: bytes | None = None,
    accent: str | None = None,
) -> bytes:
    """Render one or more Battlecards into a single deck, one slide per area."""
    heading = _hex(heading_color)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for card in cards:
        _add_card_slide(prs, blank, card, heading, None, logo, accent)
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _add_card_slide(
    prs: Presentation,
    blank: SlideLayout,
    card: Battlecard,
    navy: RGBColor,
    map_geometry: dict | None = None,
    logo: bytes | None = None,
    accent: str | None = None,
    map_image: bytes | None = None,
) -> None:
    slide = prs.slides.add_slide(blank)
    vs = card.visual_summary
    accent_rgb = _hex(accent) if accent else _GOLD
    _sidebar(slide, card, navy, map_geometry, logo, accent_rgb, map_image)
    _pillars(slide, vs.header.lifestyle_pillars)
    _messaging_columns(slide, card, navy)
    _charts_row(slide, card, navy)


# --- shape helpers ----------------------------------------------------------


def _rect(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    color: RGBColor,
) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def _box(
    slide: Slide, left: Length, top: Length, width: Length, height: Length
) -> TextFrame:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def _line(
    tf: TextFrame,
    text: str,
    *,
    size: float,
    color: RGBColor,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    first: bool = False,
    space_after: int = 2,
) -> None:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = _FONT
    run.font.color.rgb = color


def _header_bar(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    text: str,
    navy: RGBColor,
) -> None:
    _rect(slide, left, top, width, Inches(0.3), navy)
    tf = _box(
        slide, left + Inches(0.1), top + Inches(0.02), width - Inches(0.2), Inches(0.26)
    )
    _line(tf, text.upper(), size=10, color=_WHITE, bold=True, first=True)


# --- regions ----------------------------------------------------------------


def _map(slide: Slide, geometry: dict | None, accent: RGBColor) -> None:
    """Draw the catchment outline as an accent silhouette in the sidebar."""
    from .mapshape import fit_ring

    left, top = Inches(0.35), Inches(4.6)
    width, height = Inches(2.7), Inches(1.8)
    pts = fit_ring(
        geometry, int(width), int(height), pad=int(Inches(0.05)), y_down=True
    )
    if len(pts) < 3:
        return
    abs_pts = [(int(left) + int(px), int(top) + int(py)) for px, py in pts]
    try:
        builder = slide.shapes.build_freeform(abs_pts[0][0], abs_pts[0][1], scale=1.0)
        builder.add_line_segments(abs_pts[1:], close=True)
        shape = builder.convert_to_shape()
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent
        shape.line.color.rgb = _WHITE
        shape.line.width = Pt(0.75)
        shape.shadow.inherit = False
    except Exception:  # the map is decorative; never fail the export on it
        pass


def _logo(slide: Slide, logo: bytes) -> None:
    """Place the brand logo at the foot of the sidebar."""
    try:
        slide.shapes.add_picture(
            io.BytesIO(logo), Inches(0.3), Inches(6.6), height=Inches(0.6)
        )
    except Exception:  # decorative; never fail the export
        pass


def _sidebar(
    slide: Slide,
    card: Battlecard,
    navy: RGBColor,
    map_geometry: dict | None = None,
    logo: bytes | None = None,
    accent: RGBColor = _GOLD,
    map_image: bytes | None = None,
) -> None:
    vs = card.visual_summary
    h = vs.header
    stats = vs.key_statistics
    _rect(slide, 0, 0, Inches(3.4), Inches(7.5), navy)
    # Prefer the OSM basemap image for context; fall back to the vector outline.
    if map_image:
        try:
            slide.shapes.add_picture(
                io.BytesIO(map_image), Inches(0.3), Inches(4.55), width=Inches(2.8)
            )
        except Exception:  # decorative; fall back to the silhouette
            if map_geometry:
                _map(slide, map_geometry, accent)
    elif map_geometry:
        _map(slide, map_geometry, accent)

    tf = _box(slide, Inches(0.3), Inches(0.4), Inches(2.8), Inches(1.0))
    _line(tf, h.development_name.upper(), size=24, color=_WHITE, bold=True, first=True)
    location = ", ".join(p for p in [h.town, h.postcode] if p)
    if location:
        _line(tf, location, size=11, color=accent)

    tf = _box(slide, Inches(0.3), Inches(1.75), Inches(2.8), Inches(0.25))
    _line(tf, "KEY STATISTICS", size=10, color=accent, bold=True, first=True)

    tiles = [
        (f"{stats.bed_range} Bed", "HOMES"),
        (_fmt(stats.price_from, money=True), "FROM"),
        (_fmt(stats.average_household_income, money=True), "AVG HH INCOME"),
        (f"{_fmt(stats.median_age)} yrs", "MEDIAN AGE"),
        (_fmt(stats.owner_occupied_percentage, pct=True), "OWNER OCCUPIED"),
        (_fmt(stats.population_catchment), "POP. CATCHMENT"),
    ]
    for i, (value, label) in enumerate(tiles):
        col = i % 2
        row = i // 2
        left = Inches(0.3 + col * 1.5)
        top = Inches(2.15 + row * 0.85)
        tf = _box(slide, left, top, Inches(1.45), Inches(0.8))
        _line(tf, value, size=15, color=_WHITE, bold=True, first=True, space_after=0)
        _line(tf, label, size=7.5, color=_GREY)

    if logo:
        _logo(slide, logo)
    elif h.strapline:
        tf = _box(slide, Inches(0.3), Inches(6.9), Inches(2.8), Inches(0.5))
        _line(tf, h.strapline, size=10, color=_GREY, italic=True, first=True)


def _pillars(slide: Slide, pillars: list[str]) -> None:
    if not pillars:
        return
    _rect(slide, Inches(3.6), Inches(0.25), Inches(9.5), Inches(0.45), _LIGHT)
    tf = _box(slide, Inches(3.7), Inches(0.33), Inches(9.3), Inches(0.3))
    _line(
        tf,
        "   ·   ".join(pillars),
        size=12,
        color=_INK,
        italic=True,
        bold=True,
        align=PP_ALIGN.CENTER,
        first=True,
    )


def _messaging_columns(slide: Slide, card: Battlecard, navy: RGBColor) -> None:
    vs = card.visual_summary
    # Left: target audience and messaging.
    _header_bar(
        slide,
        Inches(3.6),
        Inches(0.9),
        Inches(4.6),
        "Target audience & messaging",
        navy,
    )
    tf = _box(slide, Inches(3.7), Inches(1.3), Inches(4.5), Inches(2.4))
    first = True
    lines = 0
    for m in vs.audience_messaging:
        if lines >= 7:
            break
        _line(
            tf,
            f"{m.tier.title()}: {m.audience}",
            size=10.5,
            color=_INK,
            bold=True,
            first=first,
        )
        first = False
        lines += 1
        for line in m.message_lines:
            if lines >= 7:
                break
            _line(tf, line, size=10, color=_INK)
            lines += 1

    # Right: development and location, or local amenities when AI enrichment is
    # attached (the wider-catchment card has amenities rather than site features).
    profile = card.local_area_profile
    title = "Local area & amenities" if profile else "The development & location"
    _header_bar(slide, Inches(8.45), Inches(0.9), Inches(4.65), title, navy)
    tf = _box(slide, Inches(8.55), Inches(1.3), Inches(4.55), Inches(2.4))
    if profile:
        _line(tf, profile.description, size=9.5, color=_INK, first=True)
        for cat, items in _group_amenities(profile.amenities):
            names = ", ".join(a.name for a in items[:6])
            _line(tf, f"{cat}: {names}", size=9, color=_INK)
    else:
        summary = (
            f"{vs.key_statistics.bed_range} bed homes from "
            f"{_fmt(vs.key_statistics.price_from, money=True)} in {vs.header.town}"
        )
        _line(tf, summary, size=10.5, color=_INK, bold=True, first=True)
        for feature in vs.development_features[:7]:
            _line(tf, feature, size=10, color=_INK)


def _group_amenities(amenities: list) -> list[tuple[str, list]]:
    order = ["Transport", "Retail", "Leisure", "Education", "Healthcare", "Other"]
    grouped: dict[str, list] = {}
    for a in amenities:
        grouped.setdefault(a.category, []).append(a)
    return [(c, grouped[c]) for c in order if grouped.get(c)]


def _charts_row(slide: Slide, card: Battlecard, navy: RGBColor) -> None:
    charts = card.visual_summary.charts
    top_bar = Inches(3.95)
    top_chart = Inches(4.3)
    height = Inches(2.7)

    _header_bar(slide, Inches(3.6), top_bar, Inches(3.0), "Age demographics", navy)
    _bar_chart(
        slide,
        Inches(3.6),
        top_chart,
        Inches(3.0),
        height,
        [b.label for b in charts.age_demographics],
        [_num(b.percentage) for b in charts.age_demographics],
        navy,
    )

    _header_bar(slide, Inches(6.75), top_bar, Inches(3.0), "Household income", navy)
    inc = charts.household_income
    _bar_chart(
        slide,
        Inches(6.75),
        top_chart,
        Inches(3.0),
        height,
        [
            "Mean",
            "Median",
            inc.lowest_la.name or "Lowest",
            inc.highest_la.name or "Highest",
        ],
        [
            _num(inc.mean),
            _num(inc.median),
            _num(inc.lowest_la.value),
            _num(inc.highest_la.value),
        ],
        navy,
    )

    _header_bar(slide, Inches(9.9), top_bar, Inches(3.2), "Housing tenure", navy)
    t = charts.housing_tenure
    _donut(
        slide,
        Inches(9.9),
        top_chart,
        Inches(3.2),
        height,
        ["Owns outright", "Owns with mortgage", "Social rented", "Private rented"],
        [
            _num(t.owns_outright),
            _num(t.owns_with_mortgage),
            _num(t.social_rented),
            _num(t.private_rented),
        ],
    )


def _bar_chart(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    categories: list[str],
    values: list[float],
    navy: RGBColor,
) -> None:
    data = CategoryChartData()
    data.categories = categories
    data.add_series("", values)
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data
    )
    chart = frame.chart
    chart.has_legend = False
    chart.has_title = False
    try:
        plot = chart.plots[0]
        plot.has_data_labels = False
        plot.gap_width = 60
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = navy
        chart.value_axis.has_major_gridlines = False
        chart.value_axis.visible = False
        chart.category_axis.tick_labels.font.size = Pt(8)
        chart.category_axis.tick_labels.font.name = _FONT
    except Exception:  # chart cosmetics are best effort
        pass


def _donut(
    slide: Slide,
    left: Length,
    top: Length,
    width: Length,
    height: Length,
    categories: list[str],
    values: list[float],
) -> None:
    data = CategoryChartData()
    data.categories = categories
    data.add_series("", values)
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, left, top, width, height, data
    )
    chart = frame.chart
    chart.has_title = False
    chart.has_legend = True
    try:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(8)
        chart.legend.font.name = _FONT
        points = chart.plots[0].series[0].points
        for i, point in enumerate(points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _TENURE_COLORS[i % len(_TENURE_COLORS)]
    except Exception:  # chart cosmetics are best effort
        pass
